"""Tests for all layouts — verify each renders without error and produces slides."""
import os
import pytest
from slidecraft import DeckBuilder, get_layout, list_layouts


THEME_CONFIG = {
    'font_display': 'Helvetica Neue',
    'font_body': 'Helvetica Neue',
    'colors': {
        'bg': '#FFFFFF',
        'surface': '#F5F5F5',
        'primary': '#2563EB',
    },
}


# Create a tiny test image for image zone tests
@pytest.fixture(scope='session')
def test_image(tmp_path_factory):
    """Create a minimal valid PNG for image layout tests."""
    from pptx.util import Inches
    import struct
    import zlib

    img_dir = tmp_path_factory.mktemp('images')
    img_path = str(img_dir / 'test.png')

    # Minimal 1x1 white PNG
    def make_png():
        sig = b'\x89PNG\r\n\x1a\n'
        # IHDR
        ihdr_data = struct.pack('>IIBBBBB', 1, 1, 8, 2, 0, 0, 0)
        ihdr_crc = zlib.crc32(b'IHDR' + ihdr_data) & 0xFFFFFFFF
        ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + struct.pack('>I', ihdr_crc)
        # IDAT
        raw = zlib.compress(b'\x00\xff\xff\xff')
        idat_crc = zlib.crc32(b'IDAT' + raw) & 0xFFFFFFFF
        idat = struct.pack('>I', len(raw)) + b'IDAT' + raw + struct.pack('>I', idat_crc)
        # IEND
        iend_crc = zlib.crc32(b'IEND') & 0xFFFFFFFF
        iend = struct.pack('>I', 0) + b'IEND' + struct.pack('>I', iend_crc)
        return sig + ihdr + idat + iend

    with open(img_path, 'wb') as f:
        f.write(make_png())
    return img_path


class TestLayoutRegistry:
    def test_all_21_layouts_registered(self):
        layouts = list_layouts()
        assert len(layouts) == 21

    def test_expected_layout_names(self):
        expected = {
            'title', 'section', 'content',
            'two_col_text_text', 'two_col_text_image',
            'two_col_image_text', 'two_col_image_image',
            'three_col_text_text_text', 'three_col_text_image_text',
            'three_col_image_text_text', 'three_col_text_text_image',
            'three_col_image_dtext', 'three_col_text_dimage',
            'three_col_dimage_text', 'three_col_dtext_image',
            'comparison', 'data_table', 'quote', 'profile',
            'metrics', 'closing',
        }
        assert set(list_layouts()) == expected

    def test_unknown_layout_raises(self):
        with pytest.raises(KeyError, match='Unknown layout'):
            get_layout('nonexistent')


class TestTitleLayout:
    def test_minimal(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('title')().render(db, {'layout': 'title', 'title': 'Hello'}, 1)
        assert len(db.prs.slides) == 1

    def test_full_content(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('title')().render(db, {
            'layout': 'title', 'title': 'Deck Title',
            'tagline': 'A great tagline', 'subtitle': 'By Author',
            'footer': 'Confidential',
        }, 1)
        assert len(db.prs.slides) == 1


class TestSectionLayout:
    def test_minimal(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('section')().render(db, {
            'layout': 'section', 'title': 'Part One',
        }, 1)
        assert len(db.prs.slides) == 1

    def test_with_number(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('section')().render(db, {
            'layout': 'section', 'title': 'Part One',
            'number': '01', 'subtitle': 'Introduction',
        }, 1)
        assert len(db.prs.slides) == 1


class TestContentLayout:
    def test_text_body(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('content')().render(db, {
            'layout': 'content', 'title': 'Overview',
            'body': 'This is a paragraph of text.',
        }, 1)
        assert len(db.prs.slides) == 1

    def test_bullet_list(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('content')().render(db, {
            'layout': 'content', 'title': 'Key Points',
            'body': ['Point one', 'Point two', 'Point three'],
        }, 2)
        assert len(db.prs.slides) == 1

    def test_numbered_list(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('content')().render(db, {
            'layout': 'content', 'title': 'Steps',
            'body': ['First', 'Second'], 'variant': 'numbered',
        }, 3)
        assert len(db.prs.slides) == 1


# -- Two-column layouts ------------------------------------------------

class TestTwoColTextText:
    def test_string_columns(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('two_col_text_text')().render(db, {
            'layout': 'two_col_text_text', 'title': 'Two Sides',
            'left': 'Left text', 'right': 'Right text',
        }, 1)
        assert len(db.prs.slides) == 1

    def test_dict_columns(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('two_col_text_text')().render(db, {
            'layout': 'two_col_text_text', 'title': 'Comparison',
            'variant': '60_40',
            'left': {'heading': 'Pros', 'bullets': ['Fast', 'Simple']},
            'right': {'heading': 'Cons', 'body': 'Some limitations.'},
        }, 2)
        assert len(db.prs.slides) == 1


class TestTwoColTextImage:
    def test_renders(self, test_image):
        db = DeckBuilder(THEME_CONFIG)
        db._base_dir = os.path.dirname(test_image)
        get_layout('two_col_text_image')().render(db, {
            'layout': 'two_col_text_image', 'title': 'Text + Image',
            'left': {'heading': 'Description', 'body': 'Some text.'},
            'right': os.path.basename(test_image),
        }, 1)
        assert len(db.prs.slides) == 1


class TestTwoColImageText:
    def test_renders(self, test_image):
        db = DeckBuilder(THEME_CONFIG)
        db._base_dir = os.path.dirname(test_image)
        get_layout('two_col_image_text')().render(db, {
            'layout': 'two_col_image_text', 'title': 'Image + Text',
            'left': os.path.basename(test_image),
            'right': {'heading': 'Details', 'bullets': ['Point A', 'Point B']},
        }, 1)
        assert len(db.prs.slides) == 1


class TestTwoColImageImage:
    def test_renders(self, test_image):
        db = DeckBuilder(THEME_CONFIG)
        db._base_dir = os.path.dirname(test_image)
        fname = os.path.basename(test_image)
        get_layout('two_col_image_image')().render(db, {
            'layout': 'two_col_image_image', 'title': 'Two Images',
            'left': fname, 'right': fname,
        }, 1)
        assert len(db.prs.slides) == 1


# -- Three-column equal layouts ----------------------------------------

class TestThreeColTextTextText:
    def test_renders(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('three_col_text_text_text')().render(db, {
            'layout': 'three_col_text_text_text', 'title': 'Three Text',
            'left': {'heading': 'One', 'body': 'First'},
            'center': {'heading': 'Two', 'body': 'Second'},
            'right': {'heading': 'Three', 'body': 'Third'},
        }, 1)
        assert len(db.prs.slides) == 1


class TestThreeColTextImageText:
    def test_renders(self, test_image):
        db = DeckBuilder(THEME_CONFIG)
        db._base_dir = os.path.dirname(test_image)
        get_layout('three_col_text_image_text')().render(db, {
            'layout': 'three_col_text_image_text', 'title': 'Text Image Text',
            'left': {'heading': 'Left', 'body': 'Text'},
            'center': os.path.basename(test_image),
            'right': {'heading': 'Right', 'body': 'Text'},
        }, 1)
        assert len(db.prs.slides) == 1


class TestThreeColImageTextText:
    def test_renders(self, test_image):
        db = DeckBuilder(THEME_CONFIG)
        db._base_dir = os.path.dirname(test_image)
        get_layout('three_col_image_text_text')().render(db, {
            'layout': 'three_col_image_text_text', 'title': 'Image Text Text',
            'left': os.path.basename(test_image),
            'center': {'heading': 'Mid', 'body': 'Text'},
            'right': {'heading': 'Right', 'body': 'Text'},
        }, 1)
        assert len(db.prs.slides) == 1


class TestThreeColTextTextImage:
    def test_renders(self, test_image):
        db = DeckBuilder(THEME_CONFIG)
        db._base_dir = os.path.dirname(test_image)
        get_layout('three_col_text_text_image')().render(db, {
            'layout': 'three_col_text_text_image', 'title': 'Text Text Image',
            'left': {'heading': 'Left', 'body': 'Text'},
            'center': {'heading': 'Mid', 'body': 'Text'},
            'right': os.path.basename(test_image),
        }, 1)
        assert len(db.prs.slides) == 1


# -- Three-column wide layouts -----------------------------------------

class TestThreeColImageDtext:
    def test_renders(self, test_image):
        db = DeckBuilder(THEME_CONFIG)
        db._base_dir = os.path.dirname(test_image)
        get_layout('three_col_image_dtext')().render(db, {
            'layout': 'three_col_image_dtext', 'title': 'Image + Wide Text',
            'narrow': os.path.basename(test_image),
            'wide': {'heading': 'Details', 'body': 'Longer text here.'},
        }, 1)
        assert len(db.prs.slides) == 1


class TestThreeColTextDimage:
    def test_renders(self, test_image):
        db = DeckBuilder(THEME_CONFIG)
        db._base_dir = os.path.dirname(test_image)
        get_layout('three_col_text_dimage')().render(db, {
            'layout': 'three_col_text_dimage', 'title': 'Text + Wide Image',
            'narrow': {'heading': 'Caption', 'body': 'Description.'},
            'wide': os.path.basename(test_image),
        }, 1)
        assert len(db.prs.slides) == 1


class TestThreeColDimageText:
    def test_renders(self, test_image):
        db = DeckBuilder(THEME_CONFIG)
        db._base_dir = os.path.dirname(test_image)
        get_layout('three_col_dimage_text')().render(db, {
            'layout': 'three_col_dimage_text', 'title': 'Wide Image + Text',
            'narrow': {'heading': 'Note', 'body': 'Side text.'},
            'wide': os.path.basename(test_image),
        }, 1)
        assert len(db.prs.slides) == 1


class TestThreeColDtextImage:
    def test_renders(self, test_image):
        db = DeckBuilder(THEME_CONFIG)
        db._base_dir = os.path.dirname(test_image)
        get_layout('three_col_dtext_image')().render(db, {
            'layout': 'three_col_dtext_image', 'title': 'Wide Text + Image',
            'narrow': os.path.basename(test_image),
            'wide': {'heading': 'Main Content', 'body': 'Detailed text.'},
        }, 1)
        assert len(db.prs.slides) == 1


# -- Other layouts ------------------------------------------------------

class TestComparisonLayout:
    def test_renders(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('comparison')().render(db, {
            'layout': 'comparison', 'title': 'Feature Comparison',
            'headers': ['Feature', 'Us', 'Them'],
            'rows': [['Speed', 'Fast', 'Slow'], ['Price', '$10', '$50']],
        }, 1)
        assert len(db.prs.slides) == 1

    def test_with_highlight(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('comparison')().render(db, {
            'layout': 'comparison', 'title': 'Compare',
            'headers': ['Feature', 'Us', 'Them'],
            'rows': [['Auth', 'Yes', 'No']],
            'highlight_column': 1,
        }, 2)
        assert len(db.prs.slides) == 1


class TestDataTableLayout:
    def test_simple_table(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('data_table')().render(db, {
            'layout': 'data_table', 'title': 'Revenue',
            'table': [
                ['Year', 'Revenue', 'Growth'],
                ['2024', '$1M', '50%'],
                ['2025', '$2M', '100%'],
            ],
        }, 1)
        assert len(db.prs.slides) == 1

    def test_with_sidebar(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('data_table')().render(db, {
            'layout': 'data_table', 'title': 'Financials',
            'table': [['Metric', 'Value'], ['ARR', '$1M']],
            'sidebar': {'title': 'Notes', 'items': ['Growing fast']},
        }, 2)
        assert len(db.prs.slides) == 1


class TestQuoteLayout:
    def test_minimal(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('quote')().render(db, {
            'layout': 'quote',
            'quote': 'The best software is invisible.',
            'attribution': 'Someone Famous',
        }, 1)
        assert len(db.prs.slides) == 1

    def test_with_title_and_role(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('quote')().render(db, {
            'layout': 'quote', 'title': 'Testimonial',
            'quote': 'Great product!', 'attribution': 'Jane Doe',
            'role': 'CEO at Acme',
        }, 2)
        assert len(db.prs.slides) == 1


class TestProfileLayout:
    def test_single_profile(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('profile')().render(db, {
            'layout': 'profile', 'title': 'Team',
            'profiles': [{
                'name': 'Alice', 'role': 'CEO',
                'company': 'Acme Inc.', 'bio': 'Founded 2020.',
            }],
        }, 1)
        assert len(db.prs.slides) == 1

    def test_multiple_profiles(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('profile')().render(db, {
            'layout': 'profile', 'title': 'Team',
            'profiles': [
                {'name': 'Alice', 'role': 'CEO'},
                {'name': 'Bob', 'role': 'CTO'},
            ],
        }, 2)
        assert len(db.prs.slides) == 1


class TestMetricsLayout:
    def test_single_metric(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('metrics')().render(db, {
            'layout': 'metrics', 'title': 'Traction',
            'metrics': [{'value': '10K', 'label': 'Users'}],
        }, 1)
        assert len(db.prs.slides) == 1

    def test_multiple_metrics(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('metrics')().render(db, {
            'layout': 'metrics', 'title': 'Key Numbers',
            'metrics': [
                {'value': '$1M', 'label': 'ARR'},
                {'value': '200', 'label': 'Customers'},
                {'value': '95%', 'label': 'Retention'},
            ],
        }, 2)
        assert len(db.prs.slides) == 1


class TestClosingLayout:
    def test_minimal(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('closing')().render(db, {
            'layout': 'closing', 'title': 'Thank You',
        }, 1)
        assert len(db.prs.slides) == 1

    def test_full(self):
        db = DeckBuilder(THEME_CONFIG)
        get_layout('closing')().render(db, {
            'layout': 'closing', 'title': 'The Ask',
            'headline': 'Join us',
            'bullets': ['$500K raise', '200 customers'],
            'contact': {'email': 'hello@example.com', 'url': 'example.com'},
            'closing': 'Built with care.',
        }, 10)
        assert len(db.prs.slides) == 1
