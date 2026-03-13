"""
DeckBuilder — orchestrates building a presentation from YAML content.

Usage:
    from slidecraft import DeckBuilder
    db = DeckBuilder.from_yaml('content.yaml')
    db.build()
    db.save('output.pptx')
"""
import os
import yaml
from pptx import Presentation
from pptx.util import Inches

from .theme import Theme
from .primitives import Primitives


class DeckBuilder:
    """Build a presentation from themed layouts and YAML content."""

    def __init__(self, theme_config: dict, width: float = 13.333, height: float = 7.5):
        self.theme = Theme(theme_config)
        self.prs = Presentation()
        self.prs.slide_width = Inches(width)
        self.prs.slide_height = Inches(height)
        self.p = Primitives(self.theme)
        self._blank_layout = self.prs.slide_layouts[6]
        self._content = None
        self._base_dir = '.'
        self._current_slide_data = None

    @classmethod
    def from_yaml(cls, path: str) -> 'DeckBuilder':
        """Create a DeckBuilder from a YAML content file."""
        with open(path) as f:
            content = yaml.safe_load(f)
        theme_config = content.get('theme', {})
        db = cls(theme_config)
        db._content = content
        db._base_dir = os.path.dirname(os.path.abspath(path))
        return db

    @property
    def content(self) -> dict:
        """The full parsed YAML content."""
        return self._content

    @property
    def slides_content(self) -> list:
        """The slides list from YAML content."""
        if self._content is None:
            return []
        slides = self._content.get('slides', [])
        if isinstance(slides, list):
            return slides
        # If slides is a dict (named slides), return as list of (name, data) pairs
        return list(slides.items())

    def resolve_path(self, relative_path: str) -> str:
        """Resolve a path relative to the YAML content file."""
        return os.path.join(self._base_dir, relative_path)

    def new_slide(self, bg: str = 'bg'):
        """Add a blank slide with themed background or background image."""
        slide = self.prs.slides.add_slide(self._blank_layout)

        # Resolve background image: per-slide overrides global theme
        data = self._current_slide_data or {}
        bg_image = data.get('background_image', self.theme.background_image)

        if bg_image is not None:
            resolved = self.resolve_path(bg_image)
            if not os.path.isfile(resolved):
                raise FileNotFoundError(
                    f"Background image not found: {resolved}")
            self.p.background_image(slide, resolved)
            # Apply overlay if configured (per-slide overrides global)
            overlay_cfg = data.get('overlay', self.theme.overlay)
            if overlay_cfg:
                self.p.overlay(
                    slide,
                    color=overlay_cfg.get('color', '#000000'),
                    opacity=overlay_cfg.get('opacity', 0.5))
        else:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self.theme.color(bg)

        return slide

    def build(self):
        """Build all slides from YAML content using the layout registry."""
        from .layouts import get_layout

        slides = self.slides_content
        slide_num = 0

        for slide_config in slides:
            # Support both list-of-dicts and dict-of-dicts
            if isinstance(slide_config, tuple):
                _name, data = slide_config
            else:
                data = slide_config

            layout_name = data.get('layout')
            if not layout_name:
                raise ValueError(f"Slide missing 'layout' key: {data}")

            layout_cls = get_layout(layout_name)
            slide_num += 1
            self._current_slide_data = data
            layout_cls().render(self, data, slide_num)
            self._current_slide_data = None

    def save(self, path: str):
        """Save the presentation to a file."""
        self.prs.save(path)
        print(f'Saved: {path}')
