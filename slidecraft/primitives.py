"""
Primitives — low-level PPTX building blocks.

All position/size parameters are in inches (float).
Color parameters accept theme color names (str) or RGBColor objects.
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .theme import Theme


ALIGN_MAP = {
    'left': PP_ALIGN.LEFT,
    'center': PP_ALIGN.CENTER,
    'right': PP_ALIGN.RIGHT,
}

ANCHOR_MAP = {
    'top': MSO_ANCHOR.TOP,
    'middle': MSO_ANCHOR.MIDDLE,
    'bottom': MSO_ANCHOR.BOTTOM,
}


class Primitives:
    """Low-level drawing primitives bound to a theme."""

    def __init__(self, theme: Theme):
        self.theme = theme

    def _resolve_color(self, color) -> RGBColor:
        return self.theme.color(color)

    def _resolve_align(self, align) -> PP_ALIGN:
        if isinstance(align, PP_ALIGN):
            return align
        return ALIGN_MAP.get(align, PP_ALIGN.LEFT)

    def _resolve_font(self, font=None, heading=False) -> str:
        if font is not None:
            return font
        return self.theme.font(heading)

    # ------------------------------------------------------------------
    # Text
    # ------------------------------------------------------------------

    def text_box(self, slide, left, top, width, height, text,
                 size=30, color='text', bold=False, align='left',
                 font=None, heading=False, line_spacing=1.2):
        """Add a text box. Returns the text frame."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.color.rgb = self._resolve_color(color)
        p.font.bold = bold
        p.font.name = self._resolve_font(font, heading)
        p.alignment = self._resolve_align(align)
        p.space_after = Pt(0)
        if line_spacing != 1.0:
            p.line_spacing = Pt(size * line_spacing)
        return tf

    def paragraph(self, tf, text, size=24, color='text', bold=False,
                  align='left', font=None, heading=False,
                  space_before=0, space_after=6):
        """Add a paragraph to a text frame. Returns the paragraph."""
        p = tf.add_paragraph()
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.color.rgb = self._resolve_color(color)
        p.font.bold = bold
        p.font.name = self._resolve_font(font, heading)
        p.alignment = self._resolve_align(align)
        p.space_before = Pt(space_before)
        p.space_after = Pt(space_after)
        return p

    def run(self, para, text, size=None, color=None, bold=None, font=None):
        """Add a run (inline span) to a paragraph. Returns the run."""
        r = para.add_run()
        r.text = str(text)
        if size is not None:
            r.font.size = Pt(size)
        if color is not None:
            r.font.color.rgb = self._resolve_color(color)
        if bold is not None:
            r.font.bold = bold
        if font is not None:
            r.font.name = font
        return r

    def bullets(self, slide, left, top, width, height, items,
                title=None, title_size=20, title_color='primary',
                bullet_size=18, bullet_color='text_secondary',
                heading_color=None, detail_color='text_muted',
                space_between=8):
        """Add a bullet list. Items can be strings or dicts with heading/detail.

        Returns the text frame.
        """
        if title:
            tf = self.text_box(slide, left, top, width, height,
                               title, size=title_size, color=title_color, bold=True)
        else:
            tf = self.text_box(slide, left, top, width, height,
                               '', size=bullet_size, color=bullet_color)
            tf.paragraphs[0].text = ''

        for item in items:
            if isinstance(item, str):
                self.paragraph(tf, '\u2022 ' + item,
                               size=bullet_size, color=bullet_color,
                               space_before=space_between)
            elif isinstance(item, dict):
                # Structured bullet: heading + detail
                h_color = heading_color or bullet_color
                p = self.paragraph(tf, '', size=bullet_size, color=h_color,
                                   bold=True, space_before=space_between)
                self.run(p, '\u2022 ' + item.get('heading', ''),
                         size=bullet_size, color=self._resolve_color(h_color),
                         bold=True, font=self.theme.font_body)
                if 'detail' in item:
                    self.paragraph(tf, '   ' + item['detail'],
                                   size=bullet_size - 2,
                                   color=detail_color, space_before=2)
        return tf

    # ------------------------------------------------------------------
    # Shapes
    # ------------------------------------------------------------------

    def accent_bar(self, slide, left, top, width=0.8, height=0.06, color='primary'):
        """Small colored accent bar."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._resolve_color(color)
        shape.line.fill.background()
        return shape

    def divider(self, slide, top, left=None, width=None, color='divider'):
        """Horizontal divider line."""
        if left is None:
            left = self.theme.margin
        if width is None:
            width = 13.333 - 2 * self.theme.margin
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(0.02))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._resolve_color(color)
        shape.line.fill.background()
        return shape

    def rect(self, slide, left, top, width, height,
             fill='surface', border='border', border_width=1.5, rounded=True):
        """Rounded or sharp rectangle. Returns the shape."""
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._resolve_color(fill)
        if border:
            shape.line.color.rgb = self._resolve_color(border)
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background()
        return shape

    def image(self, slide, path, left, top, width, height=None):
        """Add an image. Height defaults to width (square)."""
        if height is None:
            height = width
        return slide.shapes.add_picture(
            path, Inches(left), Inches(top), Inches(width), Inches(height))

    # ------------------------------------------------------------------
    # Tables
    # ------------------------------------------------------------------

    def table(self, slide, left, top, width, height, data,
              col_widths=None, header_color='primary',
              cell_color_fn=None, text_color_fn=None,
              font_size=16, align_fn=None):
        """Add a table from a list-of-lists.

        Optional callbacks for per-cell customization:
            cell_color_fn(row_idx, col_idx, value) -> color name or None
            text_color_fn(row_idx, col_idx, value) -> color name or None
            align_fn(row_idx, col_idx) -> 'left'|'center'|'right' or None

        Returns the table shape.
        """
        rows = len(data)
        cols = len(data[0]) if data else 0
        tbl_shape = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top), Inches(width), Inches(height))
        tbl = tbl_shape.table

        if col_widths:
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = Inches(w)

        for row_idx, row_data in enumerate(data):
            for col_idx, val in enumerate(row_data):
                cell = tbl.cell(row_idx, col_idx)
                cell.text = str(val)
                cell.fill.solid()
                p = cell.text_frame.paragraphs[0]
                p.font.name = self.theme.font_body
                p.font.size = Pt(font_size)
                cell.margin_left = Inches(0.1)
                cell.margin_top = Inches(0.04)
                cell.margin_bottom = Inches(0.04)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Default styling
                if row_idx == 0:
                    cell.fill.fore_color.rgb = self._resolve_color('surface')
                    p.font.bold = True
                    p.font.color.rgb = self._resolve_color(header_color)
                else:
                    alt = 'row_alt' if row_idx % 2 == 0 else 'bg'
                    cell.fill.fore_color.rgb = self._resolve_color(alt)
                    p.font.color.rgb = self._resolve_color('text_secondary')

                # Custom overrides
                if cell_color_fn:
                    cc = cell_color_fn(row_idx, col_idx, val)
                    if cc:
                        cell.fill.fore_color.rgb = self._resolve_color(cc)
                if text_color_fn:
                    tc = text_color_fn(row_idx, col_idx, val)
                    if tc:
                        p.font.color.rgb = self._resolve_color(tc)
                if align_fn:
                    a = align_fn(row_idx, col_idx)
                    if a:
                        p.alignment = self._resolve_align(a)

        return tbl_shape

    # ------------------------------------------------------------------
    # Common patterns
    # ------------------------------------------------------------------

    def slide_header(self, slide, num, title, subtitle=None,
                     subtitle_color='text_secondary'):
        """Accent bar + title + optional subtitle + slide number.

        Returns the y position (inches) after the header, for content placement.
        """
        m = self.theme.margin
        self.accent_bar(slide, m, self.theme.spacing('margin_top'))
        self.text_box(slide, m, self.theme.spacing('margin_top') + 0.15,
                      13.333 - 2 * m, 0.8,
                      title, size=40, color='text', bold=True, heading=True)
        self.slide_number(slide, num)
        y = 1.5
        if subtitle:
            self.text_box(slide, m, y, 13.333 - 2 * m, 0.6,
                          subtitle, size=22, color=subtitle_color)
            y = 2.2
        return y

    def slide_number(self, slide, num):
        """Slide number in bottom-right corner."""
        self.text_box(slide, 12.2, 7.0, 0.8, 0.4,
                      str(num), size=14, color='text_dim', align='right')
